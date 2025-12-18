import os
from typing import List, Dict, Optional

import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation

import io
import fitz  # PyMuPDF
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from config import DATA_DIR, IMAGES_DIR

class DocumentLoader:
    def __init__(self, data_dir: str = DATA_DIR):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]

    def _save_image(self, image_bytes, theme, filename, index):
        """辅助函数：保存图片到本地"""
        if not image_bytes:
            return None
        
        # 创建主题文件夹
        save_dir = os.path.join(IMAGES_DIR, theme)
        os.makedirs(save_dir, exist_ok=True)
        
        # 生成唯一文件名
        img_filename = f"{os.path.splitext(filename)[0]}_img_{index}.png"
        img_path = os.path.join(save_dir, img_filename)
        
        try:
            image = Image.open(io.BytesIO(image_bytes))
            # 转换为RGB防止报错
            if image.mode != "RGB":
                image = image.convert("RGB")
            image.save(img_path, format="PNG")
            return img_path
        except Exception as e:
            print(f"图片保存失败: {e}")
            return None

    def load_pdf(self, file_path: str, theme: str) -> List[Dict]:
        """提取PDF文本和图片"""
        results = []
        filename = os.path.basename(file_path)
        
        try:
            doc = fitz.open(file_path)
            for i, page in enumerate(doc):
                # 1. 提取文本
                text = page.get_text()
                if text.strip():
                    results.append({
                        "content": f"--- 第 {i+1} 页 ---\n{text}\n",
                        "image_path": None,
                        "page_number": i + 1
                    })
                
                # 2. 提取图片
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    saved_path = self._save_image(image_bytes, theme, filename, f"p{i+1}_{img_index}")
                    
                    if saved_path:
                        # 创建一个特殊的“图片块”，内容先留空，后续由Vision模型填充描述
                        results.append({
                            "content": "[IMAGE_PENDING_DESCRIPTION]", # 占位符
                            "image_path": saved_path,
                            "page_number": i + 1,
                            "is_image": True # 标记这是一个图片块
                        })
                        
        except Exception as e:
            print(f"Error reading PDF {file_path}: {e}")
        return results

    def load_pptx(self, file_path: str, theme: str) -> List[Dict]:
        """提取PPT文本和图片"""
        results = []
        filename = os.path.basename(file_path)
        
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                # 遍历形状
                for shape_idx, shape in enumerate(slide.shapes):
                    # 1. 提取文本
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                    
                    # 2. 提取图片
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_bytes = shape.image.blob
                        saved_path = self._save_image(image_bytes, theme, filename, f"s{i+1}_{shape_idx}")
                        if saved_path:
                            results.append({
                                "content": "[IMAGE_PENDING_DESCRIPTION]",
                                "image_path": saved_path,
                                "page_number": i + 1,
                                "is_image": True
                            })

                page_content = "\n".join(slide_texts)
                if page_content.strip():
                    results.append({
                        "content": f"--- 幻灯片 {i+1} ---\n{page_content}\n",
                        "image_path": None,
                        "page_number": i + 1
                    })
                    
        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}")
        return results

    # docx 和 txt 类似修改，这里省略 docx 的图片提取以节省篇幅，逻辑同上

    def load_document(self, file_path: str, theme: str = "default") -> List[Dict]:
        """入口函数"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        base_meta = {
            "filename": filename,
            "filepath": file_path,
            "filetype": ext,
        }
        
        raw_chunks = []
        if ext == ".pdf":
            raw_chunks = self.load_pdf(file_path, theme)
        elif ext == ".pptx":
            raw_chunks = self.load_pptx(file_path, theme)
        elif ext == ".txt":
            # 简单文本处理
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                raw_chunks = [{"content": f.read(), "image_path": None, "page_number": 0}]
        
        # 合并元数据
        documents = []
        for chunk in raw_chunks:
            doc = base_meta.copy()
            doc.update(chunk)
            documents.append(doc)
            
        return documents

    def load_all_documents(self, specific_dir: str = None) -> List[Dict]:
        """加载指定目录下的文档"""
        target_dir = specific_dir if specific_dir else self.data_dir
        # 获取主题名称 (文件夹名)
        theme_name = os.path.basename(target_dir) if specific_dir else "default"

        if not os.path.exists(target_dir):
            return []

        documents = []
        for root, _, files in os.walk(target_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    # 传入 theme 用于图片分类存储
                    doc_chunks = self.load_document(file_path, theme=theme_name)
                    documents.extend(doc_chunks)

        return documents